# -*- coding: utf-8 -*-
#
# Docling service wrapper using MCP with fallback to basic text extraction
#

import asyncio
from typing import Dict, Any, Optional, List
from pathlib import Path
import re


class DoclingService:
    """
    Service for parsing documents using Docling Python library.
    
    Docling supports:
    - Documents: PDF, DOCX, PPTX, XLSX, DOC, PPT, XLS
    - Web/Markup: HTML, HTM, MD, XML
    - Images: PNG, JPG, JPEG, TIFF, TIF, BMP (with OCR)
    - Audio: WAV, MP3, M4A (with transcription)
    """
    
    def __init__(self):
        self.use_docling = True  # Using Docling Python library directly
        
    async def parse_document(self, file_path: str) -> Dict[str, Any]:
        """
        Parse a document using Docling Python library.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary with parsed content, chunks, and metadata
        """
        if self.use_docling:
            return await self._parse_with_docling(file_path)
        else:
            # This fallback should never be used
            raise Exception("Docling is required - no fallback available")
    
    async def _parse_with_docling(self, file_path: str) -> Dict[str, Any]:
        """
        Parse using Docling Python library.
        Supports all Docling-compatible formats.
        """
        try:
            from docling.document_converter import DocumentConverter
            
            print(f"Docling: Converting {Path(file_path).name}")
            
            # Initialize converter
            converter = DocumentConverter()
            
            # Convert document
            result = converter.convert(file_path)
            
            # Export to markdown
            markdown_content = result.document.export_to_markdown()
            
            print(f"Docling: Extracted {len(markdown_content)} characters")
            
            # Split into chunks
            chunks = self._chunk_text(markdown_content, chunk_size=3000, overlap=200)
            
            return {
                "document_key": f"doc_{Path(file_path).stem}",
                "content": markdown_content,
                "chunks": chunks,
                "method": "docling",
                "file_type": Path(file_path).suffix
            }
        except ImportError:
            raise Exception("Docling not installed. Run: pip install docling")
        except Exception as e:
            print(f"Docling error: {e}")
            raise Exception(f"Docling parsing failed: {str(e)}")
    
    async def _parse_with_fallback(self, file_path: str, file_ext: str) -> Dict[str, Any]:
        """
        Parse using fallback libraries based on file type.
        """
        try:
            if file_ext == '.pdf':
                content = await self._parse_pdf_fallback(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                content = await self._parse_pptx_fallback(file_path)
            elif file_ext in ['.docx', '.doc']:
                content = await self._parse_docx_fallback(file_path)
            else:
                content = f"Unsupported file type: {file_ext}"
            
            # Split content into chunks
            chunks = self._chunk_text(content)
            
            return {
                "document_key": f"doc_{Path(file_path).stem}",
                "content": content,
                "chunks": chunks,
                "method": "fallback",
                "file_type": file_ext
            }
        except Exception as e:
            return {
                "document_key": f"doc_{Path(file_path).stem}",
                "content": f"Error parsing document: {str(e)}",
                "chunks": [],
                "method": "fallback",
                "error": str(e)
            }
    
    async def _parse_pdf_fallback(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF2."""
        try:
            import PyPDF2
            text_parts = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            return "\n\n".join(text_parts)
        except ImportError:
            return "PyPDF2 not installed. Install with: pip install PyPDF2"
        except Exception as e:
            return f"Error parsing PDF: {str(e)}"
    
    async def _parse_pptx_fallback(self, file_path: str) -> str:
        """Extract text from PowerPoint using python-pptx."""
        try:
            from pptx import Presentation
            text_parts = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n\n".join(text_parts)
        except ImportError:
            return "python-pptx not installed. Install with: pip install python-pptx"
        except Exception as e:
            return f"Error parsing PowerPoint: {str(e)}"
    
    async def _parse_docx_fallback(self, file_path: str) -> str:
        """Extract text from Word document using python-docx."""
        try:
            from docx import Document
            doc = Document(file_path)
            text_parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
            return "\n\n".join(text_parts)
        except ImportError:
            return "python-docx not installed. Install with: pip install python-docx"
        except Exception as e:
            return f"Error parsing Word document: {str(e)}"
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks for better context preservation.
        
        Args:
            text: Input text
            chunk_size: Target size of each chunk in characters
            overlap: Number of overlapping characters between chunks
            
        Returns:
            List of chunk dictionaries with content and metadata
        """
        if not text or len(text) < chunk_size:
            return [{"content": text, "index": 0, "length": len(text)}]
        
        chunks = []
        start = 0
        index = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundaries
            if end < len(text):
                # Look for sentence endings near the chunk boundary
                boundary_search = text[end-100:end+100]
                sentence_ends = [m.start() + end - 100 
                               for m in re.finditer(r'[.!?]\s+', boundary_search)]
                if sentence_ends:
                    # Use the closest sentence boundary
                    closest = min(sentence_ends, key=lambda x: abs(x - end))
                    end = closest
            
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "content": chunk_text,
                    "index": index,
                    "start": start,
                    "end": end,
                    "length": len(chunk_text)
                })
                index += 1
            
            # Move start forward, accounting for overlap
            start = end - overlap if end < len(text) else end
        
        return chunks

